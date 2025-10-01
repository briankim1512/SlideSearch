import os, re, hashlib, sqlite3, time, logging
import webview, win32com.client, pythoncom, win32timezone # win32timezone is needed for hidden imports
from urllib.parse import quote
from pathlib import Path
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont

PROGRAM_DATA_PATH = os.path.join(
    os.path.dirname(__file__),
    'assets',
    'data'
)
HOME_PATH = str(Path.home())
DATABASE_PATH = os.path.join(PROGRAM_DATA_PATH, 'slides.db')
LOG_PATH = os.path.join(PROGRAM_DATA_PATH, 'logs.log')

# Check if log file exists, if not create it
if not os.path.exists(LOG_PATH):
    os.makedirs(PROGRAM_DATA_PATH, exist_ok=True)
    with open(LOG_PATH, 'w'):
        pass

logging.basicConfig(
    filename=LOG_PATH,
    filemode='a',
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)

def initialize_database():
    os.makedirs(PROGRAM_DATA_PATH, exist_ok=True)
    connection = sqlite3.connect(DATABASE_PATH)
    cursor = connection.cursor()
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS slides (
            id INTEGER PRIMARY KEY,
            pptx_hash TEXT,
            slide_hash TEXT,
            pptx_name TEXT,
            pptx_modified TEXT,
            pptx_path TEXT,
            slide_number INTEGER,
            text TEXT,
            notes TEXT
        )
    ''')
    connection.commit()
    connection.close()

def get_index():
    path = os.path.join(
        os.path.dirname(__file__),
        'assets',
        'index.html'
    )
    try:
        with open(path, 'r', encoding='utf-8') as f:
            html = f.read()
    except FileNotFoundError:
        html = '<h1>index.html not found</h1>'

    return html

def md5sum(path, chunk=1<<20):
    h = hashlib.md5()
    with open(path, "rb") as f:
        for c in iter(lambda: f.read(chunk), b""):
            h.update(c)
    return h.hexdigest()

class Api:
    def __init__(self):
        self.pptx_app = None
        self.sqlite_conn = None
        self.sqlite_cursor = None

    def pick_files(self):
        logging.info('Opening file dialog...')
        return webview.windows[0].create_file_dialog(
            webview.FileDialog.OPEN,
            allow_multiple=True
        ) or []

    def export_slide_thumbnail_placeholders(self, path, slides, app):
        try:
            os.makedirs(path, exist_ok=True)
            app.Export(path, 'PNG', 800, 600)
        except Exception as e:
            logging.warning(f'Error exporting slides as images: {str(e)}')
            logging.info("Placeholder images will be used instead.")
            for idx, slide_data in enumerate(slides):
                slide_text = slide_data['text']
                img = Image.new(
                    'RGB', (800, 600),
                    color=(255, 255, 255)
                )
                d = ImageDraw.Draw(img)
                try:
                    header_font = ImageFont.truetype("malgun.ttf", 40)
                    body_font = ImageFont.truetype("malgun.ttf", 24)
                except IOError:
                    try:
                        header_font = ImageFont.truetype(
                            "arialuni.ttf",
                            40
                        )
                        body_font = ImageFont.truetype(
                            "arialuni.ttf",24
                        )
                    except IOError:
                        header_font = ImageFont.load_default()
                        body_font = ImageFont.load_default()

                margin = 10
                max_width = img.width - 2 * margin

                # Draw header "참고용 미리보기" centered
                header_text = "참고용 미리보기"
                header_bbox = header_font.getbbox(header_text)
                header_width = header_bbox[2] - header_bbox[0]
                header_height = header_bbox[3] - header_bbox[1]
                header_x = (img.width - header_width) // 2
                header_y = margin
                d.text(
                    (header_x, header_y),
                    header_text,
                    fill=(128, 128, 128),
                    font=header_font
                )

                # Prepare slide text lines
                words = slide_text[:1000].split()
                lines = []
                current_line = ""
                for word in words:
                    test_line = current_line + (
                        " " if current_line else ""
                    ) + word
                    bbox = body_font.getbbox(test_line)
                    width = bbox[2] - bbox[0]
                    if width > max_width and current_line:
                        lines.append(current_line)
                        current_line = word
                    else:
                        current_line = test_line
                if current_line:
                    lines.append(current_line)

                # Calculate total height of all lines
                line_height = body_font.getbbox('A')[3] - \
                    body_font.getbbox('A')[1] + 6
                total_text_height = len(lines) * line_height

                # Start drawing text below the header is centered
                y_start = header_y + header_height + margin
                available_height = img.height - y_start - margin
                text_y = y_start + max(
                    0, 
                    (available_height - total_text_height) // 2
                )

                for line in lines:
                    bbox = body_font.getbbox(line)
                    line_width = bbox[2] - bbox[0]
                    x = (img.width - line_width) // 2
                    d.text(
                        (x, text_y),
                        line,
                        fill=(0, 0, 0),
                        font=body_font
                    )
                    text_y += line_height
                    if text_y > img.height - margin:
                        break  # Don't draw outside the image

                img_path = os.path.join(path, f'슬라이드{idx + 1}.PNG')
                img.save(img_path)
            logging.info(f'Placeholder images saved to: {path}')

    def parse_pptx(self, path):
        logging.info(f"Parsing PPTX file: {path}")
        slides_data = []

        try:
            logging.info("Checking if file exists...")
            pptx_hash = md5sum(path)
            pptx_name = os.path.basename(path)

            self.sqlite_cursor.execute('''
                SELECT COUNT(*) FROM slides WHERE pptx_hash = ?
            ''', (pptx_hash,))
            existing_count = self.sqlite_cursor.fetchone()[0]
            if existing_count > 0:
                logging.warning(
                    f"Presentation {pptx_name} already exists in the ",
                    f"database with {existing_count} slides", 
                    ", skipping insertion."
                )
                raise Exception(
                    "Presentation already exists in the database."
                )

            prs_app = self.pptx_app.Presentations.Open(
                path,
                ReadOnly=True,
                WithWindow=False
            )
            prs_api = Presentation(path)
            pptx_modified = prs_app.BuiltInDocumentProperties(
                "Last Save Time"
            ).Value.strftime('%Y-%m-%d')

            # Check if presentation already exists in the database
            for idx, slide in enumerate(prs_api.slides):
                logging.info(f'Parsing slide: {idx + 1}')

                slide_text = '\n'.join([
                    str(shape.text).strip() for shape in sorted(
                        [s for s in slide.shapes if s.has_text_frame],
                        key=lambda s: (s.top, s.left)
                    )
                ])
                slide_text = slide_text.strip()
                slide_text = re.sub(r'\n{2,}', '\n', slide_text)
                slide_data = {
                    'pptx_name': pptx_name,
                    'pptx_path': path,
                    'pptx_modified': pptx_modified,
                    'pptx_hash': pptx_hash,
                    'slide_hash': hashlib.md5(
                        (pptx_hash + slide_text).encode()
                    ).hexdigest(),
                    'slide_number': idx + 1,
                    'text': slide_text,
                    'notes': slide.has_notes_slide and \
                        slide.notes_slide.notes_text_frame.text or '',
                }
                logging.info(f'Extracted slide data: {slide_data}')

                self.sqlite_cursor.execute('''
                    INSERT OR REPLACE INTO slides (
                        pptx_hash,
                        slide_hash, 
                        pptx_name, 
                        pptx_modified, 
                        pptx_path, 
                        slide_number, 
                        text, 
                        notes
                    )
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                ''', (
                    slide_data['pptx_hash'],
                    slide_data['slide_hash'],
                    slide_data['pptx_name'],
                    slide_data['pptx_modified'],
                    slide_data['pptx_path'],
                    slide_data['slide_number'],
                    slide_data['text'],
                    slide_data['notes']
                )) 

                slides_data.append(slide_data)
        
            slides_thumbnails_folder = os.path.join(
                PROGRAM_DATA_PATH, 
                pptx_hash
            )
            logging.info(
                "Exporting slides as images " + \
                f"to: {slides_thumbnails_folder}"
            )
            self.export_slide_thumbnail_placeholders(
                slides_thumbnails_folder, 
                slides_data, 
                prs_app
            )
        except Exception as e:
            if str(e) == "Presentation already exists in the database.":
                logging.warning(
                    f'Skipping file {pptx_name} as' + \
                    ' it already exists in the database.'
                )
            else:
                logging.error("Error parsing PPTX file:", str(e))
                logging.error("PPTX name is ", pptx_name)
                logging.error("Continuing to next file...")
                pass
        finally:
            if prs_app:
                prs_app.Close()
            self.sqlite_conn.commit()

        return slides_data

    def ingest_files(self, paths):
        try:
            logging.info(f'Ingesting PPTX files: {paths}')
            logging.info('Spooling up PowerPoint engine...')
            
            pythoncom.CoInitialize()
            self.pptx_app = win32com.client.Dispatch(
                "PowerPoint.Application"
            )
            self.sqlite_conn = sqlite3.connect(DATABASE_PATH)
            self.sqlite_cursor = self.sqlite_conn.cursor()
            results = []
            error_paths = []

            if not paths:
                return "No files selected"
            for path in paths:
                if not path.lower().endswith('.pptx'):
                    logging.warning(
                        f"File {path} is not a .pptx file, skipping."
                    )
                    continue

                try:
                    slides = self.parse_pptx(path)
                    results.extend(slides)
                except Exception as e:
                    logging.error(f"Error processing {path}: {str(e)}")
                    error_paths.append(path)
        finally:
            self.pptx_app.Quit()
            self.sqlite_conn.close()
            pythoncom.CoUninitialize()

            self.sqlite_cursor = None
            self.sqlite_conn = None
            self.pptx_app = None

            logging.info("PowerPoint engine closed.")
            if error_paths:
                logging.warning(
                    "Errors occurred with the following " + \
                    f"files: {error_paths}"
                )

        return results
    
    def search_slides(self, query):
        logging.info(f"Searching slides for query: {query}")
        
        base_query = 'SELECT pptx_name, pptx_modified, pptx_hash, ' + \
            'slide_hash, slide_number, text, notes FROM slides'
        term_query = ' AND WHERE (text LIKE ? OR notes LIKE ?)'
        title_query = ' AND WHERE pptx_name LIKE ?'
        time_range_query = ' AND WHERE pptx_modified BETWEEN ? AND ?'

        if query['text']:
            term_query = term_query.replace('?', f"'%{query['text']}%'")
            if 'WHERE' not in base_query:
                temp_query = term_query.replace('AND ', '')
                base_query += temp_query
            else:
                temp_query = term_query.replace('WHERE ', '')
                base_query += temp_query

        if query['title']:
            title_query = title_query.replace(
                '?',
                f"'%{query['title']}%'"
            )
            if 'WHERE' not in base_query:
                temp_query = title_query.replace('AND ', '')
                base_query += temp_query
            else:
                temp_query = title_query.replace('WHERE ', '')
                base_query += temp_query

        if query['time_range'][0] and query['time_range'][1]:
            time_range_query = time_range_query.replace(
                '?', 
                f"'{query['time_range'][0]}'", 
                1
            )
            time_range_query = time_range_query.replace(
                '?', 
                f"'{query['time_range'][1]}'", 
                1
            )
            if 'WHERE' not in base_query:
                temp_query = time_range_query.replace('AND ', '')
                base_query += temp_query
            else:
                temp_query = time_range_query.replace('WHERE ', '')
                base_query += temp_query

        logging.info(base_query)

        connection = sqlite3.connect(DATABASE_PATH)
        cursor = connection.cursor()
        cursor.execute(base_query)
        results = cursor.fetchall()
        connection.close()
        logging.info(f"Found {len(results)} matching slides.")

        results = [
            {
                'hash': row[3],
                'image_path': os.path.join(
                    'data',
                    row[2],
                    f'{quote("슬라이드")}{row[4]}.PNG'
                ).replace('\\', '/'),
                'pptx_name': row[0],
                'pptx_modified': row[1],
                'slide_number': row[4],
                'text': row[5],
                'notes': row[6],
            }
            for row in results
        ]

        return results
    
    def stitch_slides(self, slide_hashes):
        logging.info(f"Stitching slides with hashes: {slide_hashes}")
        if not slide_hashes:
            return "No slides selected"

        connection = sqlite3.connect(DATABASE_PATH)
        cursor = connection.cursor()
        cursor.execute(f'''
            SELECT DISTINCT pptx_path, slide_number FROM slides
            WHERE slide_hash IN ({','.join('?' for _ in slide_hashes)})
        ''', slide_hashes)
        presentations = cursor.fetchall()

        presentations_dict = {}
        for pptx_path, slide_number in presentations:
            if pptx_path not in presentations_dict:
                presentations_dict[pptx_path] = []
            presentations_dict[pptx_path].append(slide_number)

        output_folder = os.path.join(HOME_PATH, 'Downloads')
        os.makedirs(output_folder, exist_ok=True)

        try:
            logging.info('Spooling up PowerPoint engine for stitching...')
            pythoncom.CoInitialize()
            self.pptx_app = win32com.client.Dispatch(
                "PowerPoint.Application"
            )
            stitched_prs = self.pptx_app.Presentations.Add()
            source_design = None
            applied_design = False

            for pptx_path, slide_numbers in presentations_dict.items():
                logging.info(f"Processing presentation: {pptx_path}")
                source_prs = self.pptx_app.Presentations.Open(
                    pptx_path,
                    WithWindow=False
                )

                if source_design is None:
                    stitched_prs.PageSetup.SlideWidth = source_prs.\
                        PageSetup.SlideWidth
                    stitched_prs.PageSetup.SlideHeight = source_prs.\
                        PageSetup.SlideHeight 
                    source_design = source_prs.Slides(1).Design

                for slide_number in slide_numbers:
                    logging.info(f"Adding slide number: {slide_number}")
                    slide = source_prs.Slides(slide_number)
                    slide.Copy()
                    new_slide = stitched_prs.Slides.Paste()

                    if not applied_design:
                        new_slide.Design = source_design
                        applied_design = True

                source_prs.Close()

            # This code is left over as a fallback if the user wants automatic saving.
            # Currently the user is expected to manually save the stitched presentation.
            # output_name = 'stitched_' + \
            #     hashlib.md5(('-'.join(slide_hashes)).encode())\
            #         .hexdigest() + f'_{time.strftime("%Y-%m-%d")}.pptx'
            # output_path = os.path.join(output_folder, output_name)
            # stitched_prs.SaveAs(output_path)
            # stitched_prs.Close()
            # logging.info(f"Stitched presentation saved to: {output_path}")
        finally:
            del self.pptx_app
            pythoncom.CoUninitialize()
            connection.close()

        return f'Stitched presentation opened'

if __name__ == '__main__':
    logging.info("Initializing database...")
    initialize_database()
    logging.info("Starting the webview application...")
    webview.create_window(
        'My WebView App',
        'assets/index.html',
        js_api=Api(),
        width=1280,
        height=768
    )
    webview.start(http_server=True)
